"""
SK Square AI Board PPT Generator
Follows ./template/template_spec.md — A4 landscape, NanumSquare, single-column body.

Usage:
    python generate.py
    python generate.py --input input/content.md --output output/
"""
from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ------- Design tokens (sourced from template/template_spec.md) -------
BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x00, 0x00, 0x00)
PRIMARY = RGBColor(0x03, 0x47, 0x50)
SECONDARY = RGBColor(0x15, 0x60, 0x82)
ALERT = RGBColor(0xC0, 0x00, 0x00)
MUTED = RGBColor(0x66, 0x66, 0x66)

FONT = "나눔스퀘어"

SLIDE_W = Cm(29.7)
SLIDE_H = Cm(21.0)
M_TOP = Cm(1.2)
M_LEFT = Cm(2.0)
M_RIGHT = Cm(1.5)
M_BOTTOM = Cm(1.2)
DIVIDER_Y = Cm(1.5)

TITLE_TOP = Cm(1.8)
TITLE_H = Cm(1.8)
BODY_TOP = Cm(4.0)
BODY_W = SLIDE_W - M_LEFT - M_RIGHT
BODY_H = SLIDE_H - BODY_TOP - M_BOTTOM - Cm(0.8)

ROOT = Path(__file__).parent
LOGO_PATH = ROOT / "template" / "logo.png"


# ------- Data model -------
@dataclass
class Slide:
    title: str
    subtitle: str | None = None
    label: str | None = None
    bullets: list[str] = field(default_factory=list)


@dataclass
class Deck:
    title: str
    date: str
    slides: list[Slide]


# ------- Content parser -------
def parse_content(md: Path) -> Deck:
    title: str | None = None
    date = datetime.now().strftime("'%y.%m")
    slides: list[Slide] = []
    current: Slide | None = None

    for raw in md.read_text(encoding="utf-8").splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        if line.startswith("# ") and title is None:
            title = line[2:].strip()
        elif line.startswith("날짜:"):
            date = line.split(":", 1)[1].strip()
        elif line.startswith("## "):
            if current:
                slides.append(current)
            current = Slide(title=line[3:].strip())
        elif line.startswith("### ") and current is not None:
            sub = line[4:].strip()
            current.subtitle = sub
            m = re.match(r"^\[(참고|별첨)\]", sub)
            if m:
                current.label = f"[{m.group(1)}]"
        elif line.lstrip().startswith("- ") and current is not None:
            current.bullets.append(line.lstrip()[2:].strip())

    if current:
        slides.append(current)
    return Deck(title=title or "보고서", date=date, slides=slides)


# ------- Rendering helpers -------
def style(run, size: int, *, bold: bool = False, color: RGBColor = TEXT) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def add_divider(slide) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, M_LEFT, DIVIDER_Y, SLIDE_W - M_RIGHT, DIVIDER_Y
    )
    line.line.color.rgb = TEXT
    line.line.width = Pt(0.75)


def add_footer(slide, page: int, total: int) -> None:
    y = SLIDE_H - M_BOTTOM - Cm(0.4)
    w = Cm(2.2)
    x = SLIDE_W - M_RIGHT - w
    _, tf = add_textbox(slide, x, y, w, Cm(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page} / {total}"
    style(run, 9, color=TEXT)

    logo_w = Cm(1.5)
    logo_x = x - logo_w - Cm(0.2)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH), logo_x, y - Cm(0.15), height=Cm(0.8)
        )
    else:
        _, ltf = add_textbox(slide, logo_x, y, logo_w, Cm(0.5))
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.RIGHT
        lr = lp.add_run()
        lr.text = "[LOGO]"
        style(lr, 9, color=MUTED)


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_divider(slide)
    return slide


# ------- Slide builders -------
def build_cover(prs, deck: Deck, total: int) -> None:
    slide = new_slide(prs)
    _, tf = add_textbox(slide, M_LEFT, Cm(8), BODY_W, Cm(4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = deck.title
    style(run, 40, bold=True)

    _, dtf = add_textbox(slide, M_LEFT, Cm(13), BODY_W, Cm(1))
    dp = dtf.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    dr = dp.add_run()
    dr.text = deck.date
    style(dr, 16)

    add_footer(slide, 1, total)


def build_toc(prs, deck: Deck, page: int, total: int) -> None:
    slide = new_slide(prs)
    _, tf = add_textbox(slide, M_LEFT, TITLE_TOP, BODY_W, TITLE_H)
    run = tf.paragraphs[0].add_run()
    run.text = "목차"
    style(run, 24, bold=True)

    items = [s.title for s in deck.slides]
    if len(items) >= 7:
        mid = (len(items) + 1) // 2
        col_w = (BODY_W - Cm(1)) / 2
        _toc_column(slide, M_LEFT, BODY_TOP, col_w, items[:mid], 1)
        _toc_column(slide, M_LEFT + col_w + Cm(1), BODY_TOP, col_w, items[mid:], mid + 1)
    else:
        _toc_column(slide, M_LEFT, BODY_TOP, BODY_W, items, 1)

    add_footer(slide, page, total)


def _toc_column(slide, left, top, width, items, start) -> None:
    row_h = Cm(1.0)
    for i, title in enumerate(items):
        _, tf = add_textbox(slide, left, top + i * row_h, width, row_h)
        p = tf.paragraphs[0]
        n = p.add_run()
        n.text = f"{start + i:02d}. "
        style(n, 14, bold=True, color=PRIMARY)
        t = p.add_run()
        t.text = title
        style(t, 14)


def build_body(prs, s: Slide, page: int, total: int) -> None:
    slide = new_slide(prs)

    # title
    _, tf = add_textbox(slide, M_LEFT, TITLE_TOP, BODY_W, TITLE_H)
    run = tf.paragraphs[0].add_run()
    run.text = s.title
    style(run, 20, bold=True)

    # subtitle / label
    top = BODY_TOP
    if s.subtitle:
        _, stf = add_textbox(slide, M_LEFT, top, BODY_W, Cm(0.8))
        sr = stf.paragraphs[0].add_run()
        sr.text = s.subtitle
        style(sr, 14, bold=True, color=PRIMARY)
        top += Cm(1.0)

    # bullets
    if s.bullets:
        _, btf = add_textbox(slide, M_LEFT, top, BODY_W, BODY_H - (top - BODY_TOP))
        for i, b in enumerate(s.bullets):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"- {b}"
            style(run, 13)

    add_footer(slide, page, total)


# ------- Entry point -------
def generate(input_path: Path, output_dir: Path) -> Path:
    deck = parse_content(input_path)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 2 + len(deck.slides)  # cover + toc + bodies
    build_cover(prs, deck, total)
    build_toc(prs, deck, 2, total)
    for i, s in enumerate(deck.slides):
        build_body(prs, s, 3 + i, total)

    output_dir.mkdir(parents=True, exist_ok=True)
    safe = re.sub(r'[\\/:*?"<>|]', "_", deck.title).strip() or "보고서"
    out = output_dir / f"{safe}_{datetime.now():%y%m%d}.pptx"
    prs.save(out)
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description="SK Square AI Board PPT Generator")
    ap.add_argument("--input", type=Path, default=ROOT / "input" / "content.md")
    ap.add_argument("--output", type=Path, default=ROOT / "output")
    args = ap.parse_args()

    if not args.input.exists():
        print(f"[ERROR] Input not found: {args.input}", file=sys.stderr)
        return 1

    print(f"[INFO] Reading {args.input}")
    out = generate(args.input, args.output)
    print(f"[OK] Generated {out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
